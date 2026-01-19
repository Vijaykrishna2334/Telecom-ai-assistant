"""
Telecom AI Voice Assistant - Professional PowerPoint Presentation Generator
Creates a stunning demo presentation with architecture diagrams and key features
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Image paths
IMAGES = {
    "rag": r"C:\Users\Administrator\.gemini\antigravity\brain\eddf7b4c-f8e8-40a0-b6b0-49578a9da71f\uploaded_image_0_1768720412768.jpg",
    "stt": r"C:\Users\Administrator\.gemini\antigravity\brain\eddf7b4c-f8e8-40a0-b6b0-49578a9da71f\uploaded_image_1_1768720412768.png",
    "vad": r"C:\Users\Administrator\.gemini\antigravity\brain\eddf7b4c-f8e8-40a0-b6b0-49578a9da71f\uploaded_image_1_1768719882970.png",
    "llm": r"C:\Users\Administrator\.gemini\antigravity\brain\eddf7b4c-f8e8-40a0-b6b0-49578a9da71f\uploaded_image_2_1768719882970.png",
    "pipeline": r"C:\Users\Administrator\.gemini\antigravity\brain\eddf7b4c-f8e8-40a0-b6b0-49578a9da71f\uploaded_image_3_1768719882970.jpg",
    "tts": r"C:\Users\Administrator\.gemini\antigravity\brain\eddf7b4c-f8e8-40a0-b6b0-49578a9da71f\uploaded_image_4_1768719882970.png",
}

# Colors
DARK_BLUE = RGBColor(15, 23, 42)       # #0F172A
CYAN = RGBColor(6, 182, 212)           # #06B6D4
PURPLE = RGBColor(139, 92, 246)        # #8B5CF6
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(226, 232, 240)
DARK_GRAY = RGBColor(30, 41, 59)       # #1E293B

def set_slide_background(slide, color):
    """Set solid color background for slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, title, subtitle=None, title_size=44, subtitle_size=28):
    """Add title and optional subtitle to slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(subtitle_size)
        p.font.color.rgb = CYAN
        p.alignment = PP_ALIGN.LEFT

def add_bullet_points(slide, bullets, left=0.5, top=1.8, width=4.5, font_size=18):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.level = 0
        p.space_before = Pt(8)

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image to slide if it exists"""
    if os.path.exists(image_path):
        if width and height:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(image_path, Inches(left), Inches(top))
        return True
    return False

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]
    
    # ========== SLIDE 1: TITLE ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🤖 Telecom AI Voice Assistant"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Powered Customer Support for Reliance Jio"
    p.font.size = Pt(28)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Features line
    feat_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = feat_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Voice Interaction • RAG Knowledge Base • Real-Time Chat"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Date and presenter
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Capstone Project Demo | January 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 2: PROBLEM STATEMENT ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "📌 The Problem", "Customer Support Challenges in Telecom")
    
    bullets = [
        "⏱️  Long wait times - 10+ minutes on hold",
        "📞  High call volume - 1000s of repetitive queries daily",
        "💰  Expensive - ₹50-100 per call for human agents",
        "😤  Poor satisfaction - Customer frustration with delays",
        "🌙  Limited hours - Not available 24/7",
        "",
        "70% of calls are simple questions:",
        "   • \"What's the cheapest plan?\"",
        "   • \"How do I activate international roaming?\"",
        "   • \"What OTT apps are included in my plan?\""
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=22)
    
    # ========== SLIDE 3: SOLUTION OVERVIEW ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "💡 The Solution", "AI-Powered Voice Assistant")
    
    bullets = [
        "🎙️  Natural Voice Conversation - Talk like a human",
        "⚡  Instant Response - Under 500ms latency",
        "📚  Accurate Answers - From curated knowledge base",
        "🔄  Barge-in Support - Interrupt AI anytime",
        "🌐  24/7 Availability - No wait time ever",
        "🔒  100% Local - No cloud APIs, full privacy"
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=24)
    
    # Impact table
    table_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    tf = table_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Business Impact:  Wait 10min → 0sec  |  Cost ₹100 → ₹0.50  |  Hours: 8 → 24/7"
    p.font.size = Pt(20)
    p.font.color.rgb = CYAN
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 4: VOICE AI ARCHITECTURE ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🏗️ Voice AI Architecture", "9-Stage Latency War Pipeline", title_size=40)
    
    # Add pipeline image (full slide)
    add_image_safe(slide, IMAGES["pipeline"], left=0.2, top=1.2, width=9.6)
    
    # ========== SLIDE 5: SPEECH TO TEXT (WHISPER) ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🎙️ Speech to Text", "Faster-Whisper: Real-Time Voice Recognition", title_size=40)
    
    # Add STT image
    add_image_safe(slide, IMAGES["stt"], left=0.2, top=1.2, width=9.6)
    
    # ========== SLIDE 6: RAG SYSTEM ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🔍 Hybrid RAG Pipeline", "Retrieval-Augmented Generation with CRAG", title_size=40)
    
    # Add RAG image
    add_image_safe(slide, IMAGES["rag"], left=0.2, top=1.2, width=9.6)
    
    # ========== SLIDE 6: VOICE ACTIVITY DETECTION ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🎯 Voice Activity Detection", "Silero-VAD: Know When You're Speaking", title_size=40)
    
    # Add VAD image
    add_image_safe(slide, IMAGES["vad"], left=0.2, top=1.2, width=9.6)
    
    # ========== SLIDE 7: LLM ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🧠 Large Language Model", "Ollama + Llama 3.1 8B: Local Intelligence", title_size=40)
    
    # Add LLM image
    add_image_safe(slide, IMAGES["llm"], left=0.2, top=1.2, width=9.6)
    
    # ========== SLIDE 8: TEXT TO SPEECH ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🔊 Text to Speech", "Kokoro-82M: Natural Voice Output", title_size=40)
    
    # Add TTS image
    add_image_safe(slide, IMAGES["tts"], left=0.2, top=1.2, width=9.6)
    
    # ========== SLIDE 9: BARGE-IN FEATURE ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🛑 Barge-In Feature", "Interrupt the AI Anytime - Just Like Real Conversation")
    
    bullets = [
        "How It Works:",
        "  1. VAD monitors audio while bot speaks",
        "  2. Detects human speech (3+ words)",
        "  3. Echo detection filters bot's own voice",
        "  4. Immediately stops TTS playback",
        "  5. Processes new user input",
        "",
        "Technical Details:",
        "  • 15 audio chunks (~1.5s) for validation",
        "  • Whisper transcription for verification",
        "  • 70% overlap detection for echo filtering"
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=22)
    
    # ========== SLIDE 10: TECH STACK ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "⚙️ Technology Stack", "100% Open Source - No Cloud APIs Required")
    
    bullets = [
        "Backend (Python + FastAPI):",
        "  • STT: Faster-Whisper (small.en) - 4x faster transcription",
        "  • TTS: Kokoro-82M - Natural voice synthesis",
        "  • VAD: Silero-VAD - 98% accuracy, 32ms latency",
        "  • LLM: Llama 3.1 8B via Ollama - Local inference",
        "  • Vector DB: ChromaDB + BGE Embeddings",
        "  • Search: Hybrid (Vector + BM25 with RRF)",
        "",
        "Frontend (React + TypeScript):",
        "  • shadcn/ui + TailwindCSS",
        "  • WebSocket for real-time communication",
        "  • Web Audio API for voice recording"
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=20)
    
    # ========== SLIDE 11: PERFORMANCE ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "📊 Performance Metrics", "Enterprise-Grade Performance on Consumer Hardware")
    
    bullets = [
        "Hardware Requirements:",
        "  • GPU: 10+ GB VRAM (RTX 3080+ recommended)",
        "  • RAM: 32 GB",
        "  • CPU: 8+ cores",
        "  • Storage: 50 GB SSD",
        "",
        "Performance Benchmarks:",
        "  • Response Latency: < 500ms end-to-end",
        "  • STT Processing: ~800ms",
        "  • TTS Generation: ~200ms",
        "  • RAG Retrieval: ~100ms",
        "  • Concurrent Users: 100+",
        "  • Accuracy: 95%+ on telecom queries"
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=20)
    
    # ========== SLIDE 12: DEMO & CONCLUSION ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🎬 Live Demo", "Let's See It In Action!")
    
    bullets = [
        "Demo Flow:",
        "  1. Ask about budget-friendly prepaid plans",
        "  2. 🔗 https://telecom-ai-assistant.vercel.app",
        "  3. Demonstrate Barge-in interruption",
        "",
        "Key Takeaways:",
        "  ✅ Real-time natural voice conversation",
        "  ✅ Accurate knowledge retrieval (95%+)",
        "  ✅ Natural interruption support",
        "  ✅ 100% local & privacy-preserving",
        "  ✅ Scalable to enterprise deployment"
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=22)
    
    # ========== SLIDE 13: FUTURE SCOPE ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    add_title_text(slide, "🚀 Future Scope", "Expanding Capabilities")
    
    bullets = [
        "Planned Enhancements:",
        "  • Multi-language support (Hindi, Tamil, etc.)",
        "  • Emotion detection and sentiment analysis",
        "  • Integration with IVR/Call Center systems",
        "  • Voice cloning for brand-specific voices",
        "  • Mobile app with offline voice support",
        "",
        "Scalability:",
        "  • Kubernetes deployment for auto-scaling",
        "  • Load balancing across multiple GPU nodes",
        "  • Real-time analytics dashboard"
    ]
    add_bullet_points(slide, bullets, left=0.5, top=1.8, width=9, font_size=22)
    
    # ========== SLIDE 14: THANK YOU ==========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)
    
    # Thank you message
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🙏 Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.7))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(32)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GitHub: github.com/your-username/Telecom-ai-assistant"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Built with love
    love_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.5))
    tf = love_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Built with ❤️ for better customer service"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = r"d:\Running models\final capstone\Telecom-ai-assistant\Telecom_AI_Demo_v2.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
