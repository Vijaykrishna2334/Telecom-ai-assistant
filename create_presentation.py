"""
Generate PowerPoint presentation for Telecom AI Assistant
Run: python create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(15, 23, 42)
CYAN = RGBColor(6, 182, 212)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets, image_note=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(20, 30, 50)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    # Image placeholder note
    if image_note:
        img_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(4))
        tf = img_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"📷 ADD IMAGE:\n{image_note}"
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(title, left_content, right_content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(20, 30, 50)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    return slide

# ============ CREATE SLIDES ============

# Slide 1: Title
add_title_slide(
    "🎤 Voice-Powered AI Customer Support",
    "Telecom AI Assistant | Real-time Conversation | 100% Local"
)

# Slide 2: Problem
add_content_slide(
    "❌ The Problem: Customer Support Challenges",
    [
        "10+ minutes average wait time on hold",
        "1000s of repetitive queries daily",
        "₹50-100 cost per human agent call",
        "Poor customer satisfaction due to delays",
        "Limited availability - not 24/7",
        "",
        "Common questions like:",
        "  'What's the cheapest plan?'",
        "  'How do I recharge?'",
        "  'What OTT apps are included?'"
    ]
)

# Slide 3: Solution
add_two_column_slide(
    "✅ The Solution: AI Voice Assistant",
    [
        "🎙️ Natural Voice Conversation",
        "⚡ Instant Response (<500ms)",
        "📚 Accurate Knowledge Answers",
        "🔄 Barge-in Interruption Support",
        "🌐 24/7 Availability",
        "🔒 100% Local - Full Privacy"
    ],
    [
        "📊 Business Impact:",
        "",
        "Wait Time: 10min → 0 sec",
        "Cost/Query: ₹50 → ₹0.50",
        "Availability: 8hrs → 24/7",
        "Accuracy: 95%+"
    ]
)

# Slide 4: Architecture
add_content_slide(
    "🏗️ System Architecture: 9-Stage Voice Pipeline",
    [
        "1. Noise Filter → Separates voice from background",
        "2. VAD (Silero) → Detects speech start/stop",
        "3. STT (Whisper) → Converts voice to text",
        "4. Knowledge Finder (RAG) → Searches documents",
        "5. LLM (Llama 3.1) → Generates response",
        "6. TTS (Kokoro) → Converts text to voice",
        "7. Barge-in Controller → Handles interruptions",
        "",
        "⚡ All stages complete in under 1 second!"
    ],
    "Voice AI 9-Stage Architecture\n(uploaded_image_1)"
)

# Slide 5: Speech to Text
add_content_slide(
    "🎙️ Speech to Text: Faster-Whisper",
    [
        "⚡ 4x Faster than original Whisper",
        "🎯 Vocabulary Prompting for telecom terms",
        "📡 Streaming Support for real-time output",
        "🔧 CTranslate2 Optimized for GPU",
        "",
        "Configuration:",
        "  Model: small.en",
        "  Device: CUDA GPU",
        "  Beam Size: 3",
        "  Custom vocabulary: Jio, JioFiber, recharge..."
    ],
    "Faster-Whisper STT Diagram\n(uploaded_image_3)"
)

# Slide 6: VAD
add_content_slide(
    "🎚️ Voice Activity Detection: Silero-VAD",
    [
        "🎯 98% Accuracy - Distinguishes speech from noise",
        "⚡ 32ms Latency - Ultra-low delay",
        "🔒 Fully Offline - No internet required",
        "💻 CPU Efficient - Minimal resources",
        "",
        "Enables:",
        "  • Barge-in feature",
        "  • Background noise filtering",
        "  • Accurate speech boundaries"
    ],
    "Silero-VAD Diagram\n(uploaded_image_4)"
)

# Slide 7: Hybrid RAG
add_content_slide(
    "📚 Hybrid RAG: Knowledge Retrieval",
    [
        "Two Search Methods Combined:",
        "",
        "1️⃣ Vector Search (ChromaDB)",
        "   • Semantic understanding",
        "   • 'Find similar meaning'",
        "",
        "2️⃣ Keyword Search (BM25)",
        "   • Exact term matching",
        "   • 'Find exact words'",
        "",
        "🔀 RRF Fusion (k=60) combines both",
        "✅ 95%+ retrieval accuracy"
    ],
    "Hybrid RAG Pipeline\n(uploaded_image_0)"
)

# Slide 8: TTS
add_content_slide(
    "🔊 Text to Speech: Kokoro-82M",
    [
        "🎵 24kHz Audio - High fidelity output",
        "🪶 82M Parameters - Lightweight yet powerful",
        "📜 Apache 2.0 License - Open source",
        "⚡ GPU Accelerated - Fast synthesis",
        "",
        "Voice: af_heart (American Female)",
        "Natural, friendly tone - not robotic!"
    ],
    "Kokoro TTS Diagram\n(uploaded_image_2)"
)

# Slide 9: Barge-in
add_content_slide(
    "🚨 Barge-in: Interrupt the AI Anytime",
    [
        "How It Works:",
        "",
        "1. VAD monitors audio while bot speaks",
        "2. Detects human speech (3+ words)",
        "3. Echo detection filters bot's own voice",
        "4. Immediately stops TTS playback",
        "5. Processes new user input",
        "",
        "🎯 Makes conversation feel NATURAL!",
        "",
        "Demo: 'I'll interrupt the AI mid-sentence...'"
    ]
)

# Slide 10: Tech Stack
add_two_column_slide(
    "🛠️ Tech Stack: 100% Open Source",
    [
        "Backend:",
        "• Framework: FastAPI + WebSockets",
        "• STT: Faster-Whisper (small.en)",
        "• LLM: Llama 3.1 8B via Ollama",
        "• TTS: Kokoro-82M",
        "• VAD: Silero-VAD",
        "• Vector DB: ChromaDB",
        "• Embeddings: BGE-base-en-v1.5"
    ],
    [
        "Frontend:",
        "• Framework: React + TypeScript",
        "• Styling: Tailwind CSS",
        "• Audio: Web Audio API",
        "• Real-time: WebSocket",
        "",
        "🔒 No Cloud APIs!",
        "🔒 Complete Privacy!"
    ]
)

# Slide 11: Requirements
add_two_column_slide(
    "💻 System Requirements",
    [
        "Hardware:",
        "• GPU: 10+ GB VRAM (RTX 3080+)",
        "• RAM: 32 GB",
        "• CPU: 8+ cores",
        "• Storage: 50 GB SSD",
        "",
        "Current Setup:",
        "• NVIDIA RTX A5000 (24GB)",
        "• 32 GB RAM",
        "• Ubuntu Linux"
    ],
    [
        "Performance Metrics:",
        "",
        "• Response: <500ms",
        "• Concurrent Users: 100+",
        "• Uptime: 99.9%",
        "• Accuracy: 95%+",
        "",
        "VRAM Usage:",
        "• Total: ~9 GB",
        "• Whisper: 1 GB",
        "• Llama: 6 GB",
        "• Others: 2 GB"
    ]
)

# Slide 12: Demo & Conclusion
add_title_slide(
    "🎬 Live Demo Time!",
    "Let's see the Voice AI in action..."
)

# Slide 13: Thank You
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = DARK_BLUE
background.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You! Questions?"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(2))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Key Takeaways:"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "✅ Real-time voice conversation  ✅ Accurate answers  ✅ Barge-in support  ✅ 100% Local"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "Telecom_AI_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("\n📷 Don't forget to add your images to slides 4-8!")
